"""Office document processing pipeline: DOCX, PPTX, XLSX.

Each format is converted to structured text segments which feed into the
shared hierarchical chunking → embedding → Qdrant pipeline.
"""
import asyncio
import logging
from pathlib import Path
from uuid import UUID

from app.config import settings
from app.models import IngestJobMessage, ProcessingStatus
from app.services import db, embeddings, gcs, qdrant
from app.services.chunking import flush_segments as _flush_segments

logger = logging.getLogger(__name__)

# MIME → handler mapping
_MIME_TO_FORMAT: dict[str, str] = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
}


async def process_office(job: IngestJobMessage) -> dict:
    """Full Office document processing pipeline.

    Routes to the appropriate format handler (DOCX, PPTX, XLSX), extracts
    structured text with metadata, chunks, embeds, and writes to Qdrant + Postgres.

    Args:
        job: Pub/Sub message describing the upload.

    Returns:
        Dict with status, job_id, chunk_count, and processor key.
    """
    logger.info("office_processor: starting job_id=%s file=%s", job.job_id, job.filename)
    fmt = _MIME_TO_FORMAT.get(job.mime_type)
    if not fmt:
        raise ValueError(f"office_processor: unexpected mime type {job.mime_type!r}")

    await db.update_material_status(job.material_id, ProcessingStatus.PROCESSING)

    try:
        local_path = await gcs.download_file(job.gcs_path, job.job_id)

        loop = asyncio.get_running_loop()
        if fmt == "docx":
            segments = await loop.run_in_executor(
                None, _extract_docx_segments, local_path
            )
        elif fmt == "pptx":
            segments = await loop.run_in_executor(
                None, _extract_pptx_segments, local_path
            )
        else:  # xlsx
            segments = await loop.run_in_executor(
                None, _extract_xlsx_segments, local_path
            )

        logger.info("office_processor: job_id=%s extracted %d segments", job.job_id, len(segments))

        max_chars = settings.chunk_max_tokens * 4
        overlap_chars = settings.chunk_overlap_tokens * 4
        chunks = _flush_segments(segments, job.material_id, job.course_id, max_chars, overlap_chars)

        logger.info("office_processor: job_id=%s built %d chunks", job.job_id, len(chunks))

        if not chunks:
            await db.update_material_status(job.material_id, ProcessingStatus.COMPLETE, chunk_count=0)
            return {"status": "complete", "job_id": str(job.job_id), "chunk_count": 0, "processor": "office"}

        texts = [c["content"] for c in chunks]
        vectors = await embeddings.embed_texts(texts)

        chunk_ids = [c["id"] for c in chunks]
        payloads = [
            {
                "chunk_id": str(c["id"]),
                "material_id": str(c["material_id"]),
                "course_id": str(c["course_id"]),
                "content": c["content"],
                "chapter": c.get("chapter"),
                "section": c.get("section"),
                "page": c.get("page"),
                "content_type": c.get("content_type", "text"),
            }
            for c in chunks
        ]
        await qdrant.upsert_chunks(chunk_ids, vectors, payloads)
        await db.insert_chunks(chunks)
        await db.update_material_status(job.material_id, ProcessingStatus.COMPLETE, chunk_count=len(chunks))

        logger.info("office_processor: complete job_id=%s chunks=%d", job.job_id, len(chunks))
        return {
            "status": "complete",
            "job_id": str(job.job_id),
            "chunk_count": len(chunks),
            "processor": "office",
        }

    except Exception:
        logger.exception("office_processor: failed job_id=%s", job.job_id)
        await db.update_material_status(job.material_id, ProcessingStatus.FAILED)
        raise
    finally:
        try:
            local_path.unlink(missing_ok=True)
        except Exception:
            pass


# ── DOCX ──────────────────────────────────────────────────────────────────────


def _extract_docx_segments(path: Path) -> list[dict]:
    """Extract ordered text segments from a DOCX file.

    Iterates body elements in document order (paragraphs and tables interleaved)
    to preserve content ordering. Heading 1 elements set the current chapter;
    Heading 2+ set the current section. Tables are converted to Markdown.

    Args:
        path: Local path to the .docx file.

    Returns:
        List of segment dicts compatible with _flush_segments.
    """
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(str(path))
    PARA_TAG = qn("w:p")
    TABLE_TAG = qn("w:tbl")

    segments: list[dict] = []
    current_chapter: str | None = None
    current_section: str | None = None

    for child in doc.element.body:
        if child.tag == PARA_TAG:
            para = Paragraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            if style.startswith("Heading 1"):
                current_chapter = text
                current_section = None
                continue
            if style.startswith("Heading"):
                current_section = text
                continue
            segments.append(
                _seg(text, current_chapter, current_section, page=None, content_type="text")
            )

        elif child.tag == TABLE_TAG:
            table = Table(child, doc)
            md = _docx_table_to_markdown(table)
            if md:
                segments.append(
                    _seg(md, current_chapter, current_section, page=None, content_type="table")
                )

    return segments


def _docx_table_to_markdown(table) -> str:
    """Convert a python-docx Table to a Markdown table string.

    Args:
        table: A python-docx Table object.

    Returns:
        Markdown-formatted table string, or empty string if table is empty.
    """
    rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    if not rows:
        return ""

    col_count = max(len(r) for r in rows)
    # Pad short rows
    rows = [r + [""] * (col_count - len(r)) for r in rows]

    header = "| " + " | ".join(rows[0]) + " |"
    separator = "| " + " | ".join(["---"] * col_count) + " |"
    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])

    parts = [header, separator]
    if body:
        parts.append(body)
    return "\n".join(parts)


# ── PPTX ──────────────────────────────────────────────────────────────────────


def _extract_pptx_segments(path: Path) -> list[dict]:
    """Extract text segments from a PPTX file, preserving slide structure.

    Each slide produces: title segment (marks section boundary), body text
    segments, and speaker notes as a separate segment. Tables within slides
    are converted to Markdown.

    Args:
        path: Local path to the .pptx file.

    Returns:
        List of segment dicts compatible with _flush_segments.
    """
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(str(path))
    segments: list[dict] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title_text = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text.strip()

        section_label = title_text or f"Slide {slide_num}"

        # Emit slide title as section marker
        if title_text:
            segments.append(
                _seg(
                    f"[Slide {slide_num}] {title_text}",
                    chapter=None,
                    section=section_label,
                    page=slide_num,
                    content_type="text",
                )
            )

        # Body shapes (skip title, handle tables separately)
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue

            # Table shape
            if shape.has_table:
                md = _pptx_table_to_markdown(shape.table)
                if md:
                    segments.append(
                        _seg(md, chapter=None, section=section_label, page=slide_num, content_type="table")
                    )
                continue

            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    segments.append(
                        _seg(text, chapter=None, section=section_label, page=slide_num, content_type="text")
                    )

        # Speaker notes — separate chunk
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()
                if notes_text:
                    segments.append(
                        _seg(
                            f"[Speaker Notes] {notes_text}",
                            chapter=None,
                            section=section_label,
                            page=slide_num,
                            content_type="text",
                            metadata={"notes": True},
                        )
                    )

    return segments


def _pptx_table_to_markdown(table) -> str:
    """Convert a python-pptx Table to a Markdown table string.

    Args:
        table: A python-pptx Table object.

    Returns:
        Markdown-formatted table string, or empty string if table is empty.
    """
    rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    if not rows:
        return ""

    col_count = max(len(r) for r in rows)
    rows = [r + [""] * (col_count - len(r)) for r in rows]

    header = "| " + " | ".join(rows[0]) + " |"
    separator = "| " + " | ".join(["---"] * col_count) + " |"
    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])

    parts = [header, separator]
    if body:
        parts.append(body)
    return "\n".join(parts)


# ── XLSX ──────────────────────────────────────────────────────────────────────


def _extract_xlsx_segments(path: Path) -> list[dict]:
    """Extract text segments from an XLSX file.

    Each worksheet becomes a separate chunk. Non-empty sheets are converted
    to Markdown tables. Formulas are read as their cached display values.

    Args:
        path: Local path to the .xlsx file.

    Returns:
        List of segment dicts compatible with _flush_segments.
    """
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True)
    segments: list[dict] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = [
            [_cell_to_str(cell) for cell in row]
            for row in ws.iter_rows()
            if any(cell.value is not None for cell in row)
        ]
        if not rows:
            continue

        md = _rows_to_markdown(rows)
        if md:
            segments.append(
                _seg(
                    f"Sheet: {sheet_name}\n{md}",
                    chapter=sheet_name,
                    section=None,
                    page=None,
                    content_type="table",
                )
            )

    return segments


def _cell_to_str(cell) -> str:
    """Convert an openpyxl cell value to a clean string.

    Args:
        cell: An openpyxl cell object.

    Returns:
        String representation of the cell value.
    """
    if cell.value is None:
        return ""
    return str(cell.value).strip()


def _rows_to_markdown(rows: list[list[str]]) -> str:
    """Convert a list of string rows to a Markdown table.

    The first row is treated as the header.

    Args:
        rows: List of rows, each a list of cell strings.

    Returns:
        Markdown-formatted table string, or empty string if rows is empty.
    """
    if not rows:
        return ""

    col_count = max(len(r) for r in rows)
    rows = [r + [""] * (col_count - len(r)) for r in rows]

    header = "| " + " | ".join(rows[0]) + " |"
    separator = "| " + " | ".join(["---"] * col_count) + " |"
    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])

    parts = [header, separator]
    if body:
        parts.append(body)
    return "\n".join(parts)


# ── Shared helpers ─────────────────────────────────────────────────────────────


def _seg(
    text: str,
    chapter: str | None,
    section: str | None,
    page: int | None,
    content_type: str,
    metadata: dict | None = None,
) -> dict:
    """Build a segment dict for use with _flush_segments.

    Args:
        text: The text content of the segment.
        chapter: Current chapter heading, if any.
        section: Current section heading, if any.
        page: Page or slide number, if available.
        content_type: One of 'text', 'table', 'equation', etc.
        metadata: Optional extra metadata stored in the chunk.

    Returns:
        Segment dict with all required keys.
    """
    return {
        "text": text,
        "chapter": chapter,
        "section": section,
        "page": page,
        "content_type": content_type,
        "metadata": metadata or {},
    }
